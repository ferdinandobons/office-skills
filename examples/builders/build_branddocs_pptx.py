# SPDX-License-Identifier: MIT
"""Deterministic builder for the COMPLEX synthetic PPTX example template.

Produces ``examples/templates/branddocs_template.pptx``: a 100% synthetic
(made-up "BrandDocs Corp" brand, never proprietary) deck that stresses the
brand-pptx extractor/generator across as many PowerPoint component types as
python-pptx (plus a little raw lxml) can author from the default template:

  * A multi-PLACEHOLDER cover/title slide built on the default ``Title Slide``
    layout - title + subtitle + date + footer + slide-number placeholders all
    filled (the multi-placeholder cover the extractor inventories as cover
    anchors).
  * Slides on several DISTINCT layouts: ``Title Slide`` (cover),
    ``Section Header`` (agenda / section list), ``Title and Content``
    (content-text and content-table), ``Title Only`` (chart + the demo slide),
    ``Picture with Caption`` (picture), and a closing ``Section Header`` slide.
  * A real ``p14:sectionLst`` (PowerPoint *sections*) injected into
    ``presentation.xml`` via lxml - python-pptx 1.x does not model it - with
    THREE named sections ("Overview", "Financials", "Closing") whose
    ``p14:sldId`` lists point at real slides. The agenda slide's body lists
    those exact section names, so the section list and the agenda text agree.
  * A NATIVE table (``graphicFrame``/``a:tbl``) via ``shapes.add_table`` with a
    header row + body rows + a totals row.
  * A NATIVE chart (``graphicFrame``/``c:chart``) via ``shapes.add_chart`` with
    a small clustered-bar dataset (the embedded chart workbook is part of the
    package, exercising the relationship/parts walker).
  * A KPI dashboard slide with four branded stat cards and a second NATIVE line
    chart, so the example deck proves more than one chart family and more than
    one data-display layout.
  * A risk heatmap slide with a second NATIVE table whose body cells use semantic
    brand fills (navy/teal/amber/light).
  * A PICTURE: the shared, deterministic text-only BrandDocs wordmark PNG from
    ``_brandlib.branddocs_mark_png`` - no external / proprietary asset on disk -
    placed both as a free ``add_picture`` shape and as a cover logo. The
    wordmark is embedded at a 4:1 aspect ratio to keep text undistorted.
  * A logo-like GROUPED-SHAPE mark + auto-shapes tinted with the synthetic BrandDocs
    theme colors (approximating SmartArt - see NOTE below).
  * Flat, borderless brand BANDS (``_add_band``) - pure decoration carrying no
    text - on the otherwise-bare cover / agenda / closing / demo slides.
  * A DEMO / sample-content slide whose ONLY text equals a layout placeholder
    prompt (``"Click to edit Master title style"``) so the engine's
    demo-detection classifies it as a clearable demo region.

NOTE - what is APPROXIMATED / SKIPPED (python-pptx limits, by design):
  * SmartArt cannot be authored by python-pptx; it is APPROXIMATED with a
    grouped set of connected auto-shapes (a "process" row of boxes + arrows).
  * New custom slide masters / layouts cannot be created from scratch by
    python-pptx; this deck REUSES the default template's master + layouts (the
    task allows exactly this). Theme COLORS/FONTS are the ones that ship with
    the default template; a BrandDocs palette is applied at the shape level.

Reproducibility: no randomness, no wall-clock. The core-properties timestamps
are pinned to a fixed instant and the embedded PNG bytes are computed
deterministically by the shared ``_brandlib`` helper, so re-running the builder
yields a byte-identical file (CI-friendly).

Run:
    PYTHONPATH=scripts .venv/bin/python examples/builders/build_branddocs_pptx.py
"""

from __future__ import annotations

from datetime import datetime, timezone
from io import BytesIO
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor

from _brandlib import branddocs_mark_png, brand_theme_slots, freeze_ooxml
from pptx.enum.chart import (
    XL_CHART_TYPE,
    XL_DATA_LABEL_POSITION,
    XL_LABEL_POSITION,
    XL_LEGEND_POSITION,
)
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.util import Emu, Pt

OUT = Path(__file__).resolve().parents[1] / "templates" / "branddocs_template.pptx"

# ---------------------------------------------------------------------------
# Synthetic "BrandDocs Corp" brand palette (made-up; never proprietary).
# ---------------------------------------------------------------------------
BRAND_NAVY = RGBColor(0x16, 0x21, 0x3F)
BRAND_TEAL = RGBColor(0x2B, 0x7C, 0xD3)
BRAND_AMBER = RGBColor(0xE0, 0x74, 0x2B)
BRAND_LIGHT = RGBColor(0xEA, 0xF1, 0xFF)
BRAND_SLATE = RGBColor(0x57, 0x78, 0xB0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

# ---------------------------------------------------------------------------
# Full 12-slot BrandDocs clrScheme map for the in-place theme1.xml rewrite
# (PPTX-A1). The four core hexes come from the SHARED ``_brandlib`` source of
# truth (``brand_theme_slots``) so the deck palette can never drift from the
# docx/xlsx themes; the remaining slots (dk2/lt2/accent3/accent5/accent6) are
# brand-derived supporting tones. Bare 6-digit RRGGBB (no '#', no leading FF).
# ---------------------------------------------------------------------------
_SHARED_SLOTS = brand_theme_slots()
BRAND_CLRSCHEME = {
    "dk1": _SHARED_SLOTS["dk1"],  # navy   16213F (text)
    "lt1": "FFFFFF",  # white surface
    "dk2": _SHARED_SLOTS["accent1"],  # teal   2B7CD3 (secondary dark)
    "lt2": _SHARED_SLOTS["lt1"],  # light  EAF1FF
    "accent1": _SHARED_SLOTS["accent1"],  # teal   2B7CD3 (primary)
    "accent2": _SHARED_SLOTS["accent2"],  # amber  E0742B
    "accent3": "DCE7FF",  # pale blue tint
    "accent4": "5778B0",  # slate (brand-derived)
    "accent5": "1F3A6B",  # deep navy-blue (brand-derived)
    "accent6": "8FA9D8",  # soft periwinkle (brand-derived)
    "hlink": _SHARED_SLOTS["hlink"],  # teal   2B7CD3
    "folHlink": _SHARED_SLOTS["dk1"],  # navy   16213F
}

# BrandDocs theme typography pairing for the fontScheme rewrite (PPTX-A2):
# Arial (major / headings) + Calibri (minor / body).
BRAND_MAJOR_FONT = "Arial"
BRAND_MINOR_FONT = "Calibri"

# DrawingML namespace for the raw-lxml theme and transition edits.
A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"

# Default-template layout indices (verified against python-pptx 1.x default).
LO_TITLE = 0  # Title Slide      : CENTER_TITLE + SUBTITLE + date/footer/num
LO_TITLE_CONTENT = 1  # Title and Content: TITLE + OBJECT body
LO_SECTION = 2  # Section Header   : TITLE + BODY
LO_COMPARISON = 4  # Comparison       : TITLE + 2 BODY headers + 2 OBJECT bodies
LO_TITLE_ONLY = 5  # Title Only       : TITLE only
LO_PIC_CAPTION = 8  # Picture w/Caption: TITLE + PICTURE + BODY

# The exact default-template prompt string that the BODY/TITLE placeholders
# carry. A demo slide's only text must EQUAL one of these for the extractor's
# language-invariant demo-detection to fire. Captured here as a constant so the
# builder stays self-documenting; it is re-read live from the layout below to
# stay correct even if the bundled template ever changes.
TITLE_PROMPT = "Click to edit Master title style"

# Named PowerPoint sections to inject (name -> 0-based slide indices it spans).
# Filled in build() once the real slides (and their sldId values) exist.
SECTION_NS = "http://schemas.microsoft.com/office/powerpoint/2010/main"
SECTION_EXT_URI = "{521415D9-36F7-43E2-AB2F-B90AF26B5E84}"


# ---------------------------------------------------------------------------
# Small text helpers.
# ---------------------------------------------------------------------------
def _set_text(placeholder, text: str, *, size=None, bold=None, color=None) -> None:
    """Set placeholder/shape text in a single run with optional run formatting."""
    tf = placeholder.text_frame
    tf.text = text
    run = tf.paragraphs[0].runs[0]
    if size is not None:
        run.font.size = Pt(size)
    if bold is not None:
        run.font.bold = bold
    if color is not None:
        run.font.color.rgb = color


def _set_run_theme_color(run, theme_color) -> None:
    """Point a run's font color at a THEME slot (PPTX-A3).

    Uses python-pptx ``run.font.color.theme_color`` so the run inherits from the
    BrandDocs ``clrScheme`` authored in :func:`_rebrand_theme` rather than from a
    hardcoded RGB. A theme-level palette change then re-tints the surface.
    """
    run.font.color.theme_color = theme_color


def _add_theme_band(slide, left, top, width, height, theme_color):
    """Like :func:`_add_band` but the fill REFERENCES a theme slot (PPTX-A3).

    The band stays visually the same brand color (now sourced from the authored
    theme) yet resolves to a theme slot, reinforcing the deck's theme provenance.
    Carries no text, so classification is unaffected.
    """
    r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    r.fill.solid()
    r.fill.fore_color.theme_color = theme_color
    r.line.fill.background()
    r.shadow.inherit = False
    sp = r._element
    parent = sp.getparent()
    parent.remove(sp)
    parent.insert(2, sp)
    return r


def _bullets(shape, items, *, size=18, color=BRAND_NAVY) -> None:
    """Fill a shape's text frame with one bullet paragraph per (text, level)."""
    text_frame = shape.text_frame
    text_frame.clear()
    for i, (txt, level) in enumerate(items):
        para = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
        para.text = txt
        para.level = level
        run = para.runs[0]
        run.font.size = Pt(size)
        run.font.color.rgb = color


def _set_notes(slide, text: str) -> None:
    """Author a short synthetic speaker note on a slide (PPTX-A4).

    Creates the slide's ``notesSlideN.xml`` part (the notes-slide part the package
    walker inventories) and fills its text frame. Static strings only - no
    wall-clock - so the build stays byte-reproducible. The DEMO slide intentionally
    never gets notes so demo detection (its only text is the prompt) is unaffected.
    """
    slide.notes_slide.notes_text_frame.text = text


def _add_band(slide, left, top, width, height, color):
    """Add a flat, borderless, shadowless brand band pushed to the back.

    A band is a pure decoration rectangle. It carries NO text (its text frame is
    left empty), so ``_slide_texts`` skips it and it never affects the extractor's
    demo / region / cover-anchor classification - it only adds brand polish. The
    element is re-inserted right after ``nvGrpSpPr``/``grpSpPr`` (index 2) so it
    renders behind every authored shape and placeholder.
    """
    r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    r.fill.solid()
    r.fill.fore_color.rgb = color
    r.line.fill.background()  # borderless
    r.shadow.inherit = False  # flat, no default shadow
    sp = r._element  # push to back (behind authored content)
    parent = sp.getparent()
    parent.remove(sp)
    parent.insert(2, sp)
    return r


def _add_stat_card(
    slide,
    left,
    top,
    width,
    height,
    value,
    label,
    *,
    color,
    invert=False,
    theme_color=None,
):
    """Add a branded KPI card with value + label text.

    When ``theme_color`` is given the card fill REFERENCES that theme slot
    (PPTX-A3) instead of an absolute RGB, so the KPI surface inherits from the
    BrandDocs theme; ``color`` still drives the text-color choice (it should be
    the visual equivalent of the theme slot).
    """
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    box.fill.solid()
    if theme_color is not None:
        box.fill.fore_color.theme_color = theme_color
    else:
        box.fill.fore_color.rgb = color
    box.line.fill.background()
    box.shadow.inherit = False
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Emu(137160)
    tf.margin_right = Emu(137160)
    tf.margin_top = Emu(91440)
    tf.margin_bottom = Emu(91440)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = value
    r = p.runs[0]
    r.font.size = Pt(24)
    r.font.bold = True
    r.font.color.rgb = WHITE if invert else BRAND_NAVY
    lp = tf.add_paragraph()
    lp.text = label
    lr = lp.runs[0]
    lr.font.size = Pt(10.5)
    lr.font.bold = True
    lr.font.color.rgb = WHITE if invert else BRAND_SLATE
    return box


# ---------------------------------------------------------------------------
# Slide builders. Each returns the created slide so build() can collect sldIds.
# ---------------------------------------------------------------------------
def _add_cover(prs, png_path: Path):
    """Multi-PLACEHOLDER cover: title + subtitle + date + footer + slide-number.

    Every placeholder the ``Title Slide`` layout exposes is filled, so the
    extractor's cover-anchor inventory sees a genuine multi-placeholder cover.
    A small 'logo' picture is dropped in the top-left corner.
    """
    layout = prs.slide_layouts[LO_TITLE]
    slide = prs.slides.add_slide(layout)
    sw = prs.slide_width
    # Brand bands: a navy header band + a thin amber accent under it (decoration
    # only, no text -> classification unaffected). (P2) The cover header/accent
    # bands REFERENCE the DARK_1 / ACCENT_2 theme slots (PPTX-A3) so a theme-level
    # palette change re-tints the cover; the visual brand colors are unchanged.
    _add_theme_band(slide, Emu(0), Emu(0), sw, Emu(228600), MSO_THEME_COLOR.DARK_1)
    _add_theme_band(
        slide, Emu(0), Emu(228600), sw, Emu(45720), MSO_THEME_COLOR.ACCENT_2
    )
    _add_theme_band(
        slide,
        Emu(6629400),
        Emu(0),
        Emu(2514600),
        prs.slide_height,
        MSO_THEME_COLOR.DARK_1,
    )
    _add_theme_band(
        slide,
        Emu(6537960),
        Emu(0),
        Emu(91440),
        prs.slide_height,
        MSO_THEME_COLOR.ACCENT_2,
    )
    ph = {p.placeholder_format.idx: p for p in slide.placeholders}
    # Keep ALL 5 placeholders filled (cover_anchors >= 5); nudge the title/subtitle
    # up into a tidier stack below the header band. (P7)
    _set_text(
        ph[0],
        "BrandDocs Corp Quarterly Business Review",
        size=40,
        bold=True,
        color=BRAND_NAVY,
    )
    # Point the cover title run at the DARK_1 theme slot so it inherits from the
    # BrandDocs theme authored in _rebrand_theme (PPTX-A3); the visual color is
    # unchanged (dk1 == navy) but the surface now resolves to a theme slot.
    _set_run_theme_color(ph[0].text_frame.paragraphs[0].runs[0], MSO_THEME_COLOR.DARK_1)
    ph[0].left, ph[0].top, ph[0].width = Emu(182880), Emu(2057400), Emu(5943600)
    _set_text(
        ph[1], "FY2026 - Performance, Outlook & Initiatives", size=20, color=BRAND_TEAL
    )
    ph[1].left, ph[1].top, ph[1].width = Emu(685800), Emu(3429000), Emu(5486400)
    if 10 in ph:  # DATE placeholder
        _set_text(ph[10], "January 15, 2026", size=12, color=BRAND_SLATE)
    if 11 in ph:  # FOOTER placeholder
        _set_text(
            ph[11],
            "BrandDocs Corp - Confidential (synthetic sample)",
            size=12,
            color=BRAND_SLATE,
        )
    if 12 in ph:  # SLIDE_NUMBER placeholder
        _set_text(ph[12], "1", size=12, color=BRAND_SLATE)
    # Thin amber rule under the title. (P7)
    _add_band(slide, Emu(685800), Emu(3200400), Emu(2286000), Emu(45720), BRAND_AMBER)
    # Logo picture in the corner: the shared generated BrandDocs wordmark.
    slide.shapes.add_picture(
        str(png_path), Emu(365760), Emu(548640), width=Emu(1645920), height=Emu(411480)
    )
    # The cover KPI cards reference theme slots (PPTX-A3): teal -> ACCENT_1,
    # amber -> ACCENT_2 (the middle card stays a light absolute fill).
    _add_stat_card(
        slide,
        Emu(6903720),
        Emu(1371600),
        Emu(1943100),
        Emu(823000),
        "18%",
        "YoY revenue growth",
        color=BRAND_TEAL,
        invert=True,
        theme_color=MSO_THEME_COLOR.ACCENT_1,
    )
    _add_stat_card(
        slide,
        Emu(6903720),
        Emu(2453640),
        Emu(1943100),
        Emu(823000),
        "92",
        "Brand health index",
        color=BRAND_LIGHT,
    )
    _add_stat_card(
        slide,
        Emu(6903720),
        Emu(3535680),
        Emu(1943100),
        Emu(823000),
        "3",
        "Office formats audited",
        color=BRAND_AMBER,
        invert=True,
        theme_color=MSO_THEME_COLOR.ACCENT_2,
    )
    return slide


def _add_agenda(prs, section_names):
    """Agenda / section-list slide: body lists the deck's real section names.

    Built on the ``Section Header`` layout (TITLE + BODY). The body text is the
    exact list of ``p14:section`` names injected later, so the agenda and the
    real section list agree (the extractor surfaces both).
    """
    layout = prs.slide_layouts[LO_SECTION]
    slide = prs.slides.add_slide(layout)
    ph = {p.placeholder_format.idx: p for p in slide.placeholders}
    # The Section Header layout ships TITLE *below* BODY (title top=4.82in,
    # body top=3.18in), which overlaps. Re-stack the SAME two placeholders into a
    # clean title -> body order (override the layout coords). (P1)
    ph[0].left, ph[0].top, ph[0].width, ph[0].height = (
        Emu(685800),
        Emu(2057400),
        Emu(7772400),
        Emu(1143000),
    )
    ph[1].left, ph[1].top, ph[1].width, ph[1].height = (
        Emu(685800),
        Emu(3429000),
        Emu(7772400),
        Emu(2286000),
    )
    # Left amber accent bar (decoration only, no text). (P2) References the
    # ACCENT_2 theme slot so the agenda bar inherits from the brand theme (A3).
    _add_theme_band(
        slide,
        Emu(457200),
        Emu(2057400),
        Emu(91440),
        Emu(3429000),
        MSO_THEME_COLOR.ACCENT_2,
    )
    _set_text(ph[0], "Agenda", size=36, bold=True, color=BRAND_NAVY)
    # Anchor the body to the TOP of its frame so the list sits right below the
    # title instead of dropping to the bottom of the tall body box. (residual fix)
    ph[1].text_frame.vertical_anchor = MSO_ANCHOR.TOP
    _bullets(
        ph[1],
        [(f"{i + 1}. {name}", 0) for i, name in enumerate(section_names)],
        size=20,
    )
    _set_notes(
        slide,
        "Speaker notes: walk the three sections top to bottom; flag that the "
        "Financials block carries the revenue table, charts and the mix doughnut.",
    )
    return slide


def _add_content_text(prs):
    """Content-text slide: title + multi-level bulleted body (Title and Content)."""
    layout = prs.slide_layouts[LO_TITLE_CONTENT]
    slide = prs.slides.add_slide(layout)
    ph = {p.placeholder_format.idx: p for p in slide.placeholders}
    _set_text(ph[0], "Executive Summary", size=32, bold=True, color=BRAND_NAVY)
    _bullets(
        ph[1],
        [
            (
                "BrandDocs Corp grew net revenue 18% YoY across all three reported regions.",
                0,
            ),
            ("Gross margin expanded 230 bps on supply-chain efficiencies.", 0),
            ("North region led growth; East rebounded after Q1 softness.", 1),
            ("FY2026 outlook raised on a stronger services pipeline.", 0),
            ("Key risk: input-cost volatility in the hardware segment.", 1),
        ],
        size=18,
    )
    # Speaker notes: a standard deck artifact (the notes slide + its text frame).
    slide.notes_slide.notes_text_frame.text = (
        "Speaker notes: lead with the 18% YoY headline, then the 230 bps margin "
        "expansion. Pause on the input-cost risk before the regional breakdown."
    )
    return slide


def _add_kpi_dashboard(prs):
    """Board-style KPI dashboard: stat cards + native line chart."""
    layout = prs.slide_layouts[LO_TITLE_ONLY]
    slide = prs.slides.add_slide(layout)
    ph = {p.placeholder_format.idx: p for p in slide.placeholders}
    _set_text(ph[0], "Performance Snapshot", size=30, bold=True, color=BRAND_NAVY)

    cards = [
        ("$14.2M", "FY net revenue", BRAND_NAVY, True),
        ("+230 bps", "Gross margin lift", BRAND_TEAL, True),
        ("86%", "Workflow adoption", BRAND_LIGHT, False),
        ("4", "Priority risks", BRAND_AMBER, True),
    ]
    left = Emu(457200)
    for i, (value, label, color, invert) in enumerate(cards):
        _add_stat_card(
            slide,
            Emu(int(left) + i * 2171700),
            Emu(1371600),
            Emu(1943100),
            Emu(1005840),
            value,
            label,
            color=color,
            invert=invert,
        )

    chart_data = CategoryChartData()
    chart_data.categories = ["Jan", "Feb", "Mar", "Apr", "May", "Jun"]
    chart_data.add_series("Brand health", (74, 78, 81, 84, 88, 92))
    chart_data.add_series("Adoption", (61, 66, 70, 76, 81, 86))
    gframe = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE_MARKERS,
        Emu(685800),
        Emu(2926080),
        Emu(5486400),
        Emu(2743200),
        chart_data,
    )
    chart = gframe.chart
    chart.has_title = True
    chart.chart_title.text_frame.text = "Health and adoption trend"
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    # Brand the chart-level font to the BrandDocs major font (PPTX-A10): legend +
    # axis labels render in Arial, matching the theme heading pairing.
    chart.font.name = BRAND_MAJOR_FONT
    chart.font.size = Pt(11)
    for ser, col in zip(chart.plots[0].series, (BRAND_TEAL, BRAND_AMBER)):
        ser.format.line.color.rgb = col
        ser.format.line.width = Pt(2.25)
    # Data labels + a value-axis number format (PPTX-A10): the chart part now
    # carries c:dLbls / c:numFmt, board-ready styling read by the chart inventory.
    plot = chart.plots[0]
    plot.has_data_labels = True
    data_labels = plot.data_labels
    data_labels.number_format = "0"
    data_labels.number_format_is_linked = False
    data_labels.position = XL_LABEL_POSITION.ABOVE
    data_labels.font.size = Pt(9)
    data_labels.font.color.rgb = BRAND_NAVY
    value_axis = chart.value_axis
    value_axis.tick_labels.number_format = "0"
    value_axis.tick_labels.number_format_is_linked = False
    value_axis.tick_labels.font.color.rgb = BRAND_NAVY
    call = slide.shapes.add_textbox(
        Emu(6537960), Emu(3017520), Emu(2286000), Emu(2286000)
    )
    tf = call.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(137160)
    p = tf.paragraphs[0]
    p.text = "Primary signal"
    p.runs[0].font.bold = True
    p.runs[0].font.size = Pt(16)
    p.runs[0].font.color.rgb = BRAND_NAVY
    bp = tf.add_paragraph()
    bp.text = "Outputs now carry visual proof artifacts, moving QA from spot-checking to repeatable evidence."
    bp.runs[0].font.size = Pt(13)
    bp.runs[0].font.color.rgb = BRAND_SLATE
    _set_notes(
        slide,
        "Speaker notes: anchor on the four headline KPIs, then read the line "
        "chart left to right - brand health and adoption both trend up through June.",
    )
    return slide


def _add_content_table(prs):
    """Content-table slide: title + a NATIVE table (graphicFrame/a:tbl).

    A 6-row x 5-col table with a MERGED title band spanning all columns (PPTX-A7,
    ``a:tc`` gridSpan via ``cell.merge``), a styled header row, body rows tinted
    with the BrandDocs palette, and a totals row - exercises the table extractor
    path including merged-cell geometry.
    """
    layout = prs.slide_layouts[LO_TITLE_ONLY]
    slide = prs.slides.add_slide(layout)
    ph = {p.placeholder_format.idx: p for p in slide.placeholders}
    _set_text(
        ph[0], "Regional Revenue (USD thousands)", size=28, bold=True, color=BRAND_NAVY
    )

    rows, cols = 6, 5
    # Table spans 1.8-5.0in vertically: top=1.8in, height=3.2in, leaving a clean
    # margin above the 7.5in slide bottom. (P8)
    left, top, width, height = Emu(457200), Emu(1645920), Emu(8229600), Emu(2926080)
    gtable = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = gtable.table
    headers = ["Region", "Q1", "Q2", "Q3", "Q4"]
    body = [
        ["North", "320", "351", "372", "410"],
        ["South", "210", "228", "241", "265"],
        ["East", "175", "168", "199", "232"],
    ]
    totals = ["Total", "705", "747", "812", "907"]

    # Row 0: a MERGED title band spanning all five columns (PPTX-A7). cell.merge
    # spans the origin cell across the row (a:tc gridSpan), giving the table a
    # polished navy/white title row.
    band = table.cell(0, 0)
    band.merge(table.cell(0, cols - 1))
    band.text = "BrandDocs Corp - FY2026 (USD thousands)"
    band.fill.solid()
    band.fill.fore_color.rgb = BRAND_NAVY
    band_run = band.text_frame.paragraphs[0].runs[0]
    band_run.font.bold = True
    band_run.font.size = Pt(15)
    band_run.font.color.rgb = WHITE
    band.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Row 1: column headers.
    for c, label in enumerate(headers):
        cell = table.cell(1, c)
        cell.text = label
        cell.fill.solid()
        cell.fill.fore_color.rgb = BRAND_NAVY
        run = cell.text_frame.paragraphs[0].runs[0]
        run.font.bold = True
        run.font.size = Pt(16)
        run.font.color.rgb = WHITE
    # Rows 2-4: body, alternating tint.
    for r, datarow in enumerate(body, start=2):
        for c, val in enumerate(datarow):
            cell = table.cell(r, c)
            cell.text = val
            cell.fill.solid()
            cell.fill.fore_color.rgb = BRAND_LIGHT if r % 2 else WHITE
            run = cell.text_frame.paragraphs[0].runs[0]
            run.font.size = Pt(14)
            run.font.color.rgb = BRAND_NAVY
    # Row 5: totals.
    for c, val in enumerate(totals):
        cell = table.cell(rows - 1, c)
        cell.text = val
        cell.fill.solid()
        cell.fill.fore_color.rgb = BRAND_TEAL
        run = cell.text_frame.paragraphs[0].runs[0]
        run.font.bold = True
        run.font.size = Pt(15)
        run.font.color.rgb = WHITE
    _set_notes(
        slide,
        "Speaker notes: the merged title band names the entity and units; read "
        "the totals row last - Q4 at 907 is the strongest quarter.",
    )
    return slide


def _add_risk_heatmap(prs):
    """Native table used as a compact risk/readiness heatmap."""
    layout = prs.slide_layouts[LO_TITLE_ONLY]
    slide = prs.slides.add_slide(layout)
    ph = {p.placeholder_format.idx: p for p in slide.placeholders}
    _set_text(ph[0], "Risk & Readiness Heatmap", size=28, bold=True, color=BRAND_NAVY)
    rows, cols = 5, 4
    gtable = slide.shapes.add_table(
        rows, cols, Emu(685800), Emu(1600200), Emu(7772400), Emu(3657600)
    )
    table = gtable.table
    headers = ["Area", "Signal", "Readiness", "Action"]
    body = [
        ["Renderer", "Sandbox drift", "Amber", "Run doctor before deep QA"],
        ["Template", "Stale field cache", "Teal", "Refresh generated indexes"],
        ["Workbook", "Formula loss", "Navy", "Diff formulas shell vs output"],
        ["Deck", "Native object loss", "Amber", "Review component warnings"],
    ]
    for c, label in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = label
        cell.fill.solid()
        cell.fill.fore_color.rgb = BRAND_NAVY
        run = cell.text_frame.paragraphs[0].runs[0]
        run.font.bold = True
        run.font.color.rgb = WHITE
        run.font.size = Pt(13)
    color_for = {"Amber": BRAND_AMBER, "Teal": BRAND_TEAL, "Navy": BRAND_NAVY}
    for r, row in enumerate(body, start=1):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = val
            cell.fill.solid()
            cell.fill.fore_color.rgb = BRAND_LIGHT
            if c == 2:
                cell.fill.fore_color.rgb = color_for[val]
            run = cell.text_frame.paragraphs[0].runs[0]
            run.font.size = Pt(12)
            run.font.color.rgb = WHITE if c == 2 else BRAND_NAVY
            run.font.bold = c == 2
    _set_notes(
        slide,
        "Speaker notes: the Readiness column colour-codes each area; call out the "
        "two amber rows (renderer, deck) as the items needing an owner this week.",
    )
    return slide


def _add_chart(prs):
    """Chart slide: title + a NATIVE clustered-bar chart (graphicFrame/c:chart)."""
    layout = prs.slide_layouts[LO_TITLE_ONLY]
    slide = prs.slides.add_slide(layout)
    ph = {p.placeholder_format.idx: p for p in slide.placeholders}
    _set_text(
        ph[0], "Quarterly Net Revenue by Region", size=28, bold=True, color=BRAND_NAVY
    )

    chart_data = CategoryChartData()
    chart_data.categories = ["Q1", "Q2", "Q3", "Q4"]
    chart_data.add_series("North", (320, 351, 372, 410))
    chart_data.add_series("South", (210, 228, 241, 265))
    chart_data.add_series("East", (175, 168, 199, 232))

    left, top, width, height = Emu(457200), Emu(1828800), Emu(8229600), Emu(4114800)
    gframe = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, left, top, width, height, chart_data
    )
    chart = gframe.chart
    chart.has_title = True
    chart.chart_title.text_frame.text = "BrandDocs Corp - FY2026 Revenue"
    # Brand the title font. (P3)
    title_font = chart.chart_title.text_frame.paragraphs[0].font
    title_font.size = Pt(16)
    title_font.bold = True
    title_font.color.rgb = BRAND_NAVY
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    # Brand the 3 series in navy/teal/amber (stays a native 3-series chart). (P3)
    for ser, col in zip(chart.plots[0].series, (BRAND_NAVY, BRAND_TEAL, BRAND_AMBER)):
        ser.format.fill.solid()
        ser.format.fill.fore_color.rgb = col
    # Soften the value-axis gridlines and brand the axis tick labels. (P3)
    value_axis = chart.value_axis
    value_axis.has_major_gridlines = True
    value_axis.major_gridlines.format.line.color.rgb = BRAND_LIGHT
    value_axis.tick_labels.font.color.rgb = BRAND_NAVY
    value_axis.tick_labels.number_format = "0"
    value_axis.tick_labels.number_format_is_linked = False
    chart.category_axis.tick_labels.font.color.rgb = BRAND_NAVY
    # Data labels on the bar chart too (PPTX-A6 optional): the chart part gains
    # c:dLbls so each column shows its value, outside-end, in navy.
    plot = chart.plots[0]
    plot.has_data_labels = True
    plot.data_labels.number_format = "0"
    plot.data_labels.number_format_is_linked = False
    plot.data_labels.position = XL_DATA_LABEL_POSITION.OUTSIDE_END
    plot.data_labels.font.size = Pt(9)
    plot.data_labels.font.color.rgb = BRAND_NAVY
    _set_notes(
        slide,
        "Speaker notes: North leads every quarter; East is the recovery story "
        "after the Q1 dip - point at the rebound from 168 to 232.",
    )
    return slide


def _add_revenue_mix(prs):
    """Financials doughnut chart (PPTX-A6): a third chart FAMILY on Title Only.

    A native ``XL_CHART_TYPE.DOUGHNUT`` of the FY revenue mix by region with
    per-point brand fills (navy/teal/amber/slate), percentage data labels, and a
    branded title - adds a chart family the deck did not previously carry plus a
    second embedded chart workbook part. Static data, no wall-clock.
    """
    layout = prs.slide_layouts[LO_TITLE_ONLY]
    slide = prs.slides.add_slide(layout)
    ph = {p.placeholder_format.idx: p for p in slide.placeholders}
    _set_text(
        ph[0], "FY2026 Revenue Mix by Region", size=28, bold=True, color=BRAND_NAVY
    )

    chart_data = CategoryChartData()
    chart_data.categories = ["North", "South", "East", "West"]
    # Shares sum to 100% (kept as fractions so the 0% number format reads cleanly).
    chart_data.add_series("Revenue mix", (0.41, 0.27, 0.20, 0.12))

    left, top, width, height = Emu(1828800), Emu(1828800), Emu(5486400), Emu(4114800)
    gframe = slide.shapes.add_chart(
        XL_CHART_TYPE.DOUGHNUT, left, top, width, height, chart_data
    )
    chart = gframe.chart
    chart.has_title = True
    chart.chart_title.text_frame.text = "Revenue mix by region"
    title_font = chart.chart_title.text_frame.paragraphs[0].font
    title_font.size = Pt(16)
    title_font.bold = True
    title_font.color.rgb = BRAND_NAVY
    chart.font.name = BRAND_MAJOR_FONT
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.RIGHT
    chart.legend.include_in_layout = False
    # Per-point brand fills (navy/teal/amber/slate) on the single doughnut series.
    points = chart.plots[0].series[0].points
    for point, col in zip(points, (BRAND_NAVY, BRAND_TEAL, BRAND_AMBER, BRAND_SLATE)):
        point.format.fill.solid()
        point.format.fill.fore_color.rgb = col
    # Percentage data labels (PPTX-A6): c:dLbls + c:numFmt '0%'.
    plot = chart.plots[0]
    plot.has_data_labels = True
    plot.data_labels.number_format = "0%"
    plot.data_labels.number_format_is_linked = False
    plot.data_labels.position = XL_LABEL_POSITION.CENTER
    plot.data_labels.font.size = Pt(11)
    plot.data_labels.font.bold = True
    plot.data_labels.font.color.rgb = WHITE
    return slide


def _add_picture_slide(prs, png_path: Path):
    """Picture slide on ``Picture with Caption``: TITLE + PICTURE ph + caption."""
    layout = prs.slide_layouts[LO_PIC_CAPTION]
    slide = prs.slides.add_slide(layout)
    ph = {p.placeholder_format.idx: p for p in slide.placeholders}
    # Fill the PICTURE placeholder (idx 1) with the shared text-only wordmark.
    # The layout's picture box is much taller than a logo, so reshape the result
    # to a wide 4:1 frame after insertion.
    if 1 in ph:
        try:
            pic = ph[1].insert_picture(str(png_path))
            # Zero the crop python-pptx applies when fitting -> undistorted.
            pic.crop_left = pic.crop_right = pic.crop_top = pic.crop_bottom = 0
            pic.left, pic.top, pic.width, pic.height = (
                Emu(1828800),
                Emu(1371600),
                Emu(4754880),
                Emu(1188720),
            )
        except Exception:
            slide.shapes.add_picture(
                str(png_path),
                Emu(1828800),
                Emu(1371600),
                width=Emu(4754880),
                height=Emu(1188720),
            )
    else:
        slide.shapes.add_picture(
            str(png_path),
            Emu(1828800),
            Emu(1371600),
            width=Emu(4754880),
            height=Emu(1188720),
        )
    # Reposition the TITLE and caption BODY into a tidy stack below the picture
    # box (bottom ~3.5in). All three placeholders stay present/filled. (P12)
    _set_text(ph[0], "The BrandDocs Wordmark", size=28, bold=True, color=BRAND_NAVY)
    ph[0].top = Emu(3429000)
    if 2 in ph:
        ph[2].top = Emu(4114800)
        _bullets(
            ph[2],
            [
                ("A synthetic, generated logo - no proprietary asset.", 0),
                ("The logo is simply the BrandDocs name, with no icon or badge.", 0),
                ("Brand guidelines:", 0),
            ],
            size=16,
        )
        # External run-level hyperlink on the caption (PPTX-A5): a real
        # a:hlinkClick + slide relationship to the brand-guidelines URL. The
        # hyperlink run inherits the hlink theme slot authored in _rebrand_theme.
        link_para = ph[2].text_frame.paragraphs[-1]
        link_run = link_para.add_run()
        link_run.text = " example.com/brand"
        link_run.font.size = Pt(16)
        link_run.hyperlink.address = "https://example.com/brand"
    # A second, free-floating copy via add_picture to exercise that path too.
    slide.shapes.add_picture(
        str(png_path), Emu(457200), Emu(548640), width=Emu(1371600), height=Emu(342900)
    )
    _set_notes(
        slide,
        "Speaker notes: the wordmark is generated, not a proprietary asset; the "
        "caption links to the brand guidelines for the full usage rules.",
    )
    return slide


def _add_smartart_approx(prs):
    """Closing/SmartArt-approximation slide: a grouped 'process' of shapes.

    python-pptx CANNOT author SmartArt, so this APPROXIMATES a 3-step process
    diagram with grouped rounded rectangles connected by arrows, tinted with the
    BrandDocs palette. Documented as an approximation in the module docstring.
    """
    layout = prs.slide_layouts[LO_TITLE_ONLY]
    slide = prs.slides.add_slide(layout)
    ph = {p.placeholder_format.idx: p for p in slide.placeholders}
    _set_text(
        ph[0],
        "Our Approach (SmartArt-style process)",
        size=28,
        bold=True,
        color=BRAND_NAVY,
    )

    steps = [
        ("Discover", BRAND_NAVY, "Validate the opportunity"),
        ("Build", BRAND_TEAL, "Ship the MVP"),
        ("Scale", BRAND_AMBER, "Grow to new regions"),
    ]
    box_w, box_h = Emu(2286000), Emu(1143000)
    top = Emu(2743200)
    gap = Emu(457200)
    x = Emu(685800)
    centers_y = int(top) + int(box_h) // 2
    prev_right = None
    for label, color, sub in steps:
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, top, box_w, box_h)
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.line.color.rgb = BRAND_SLATE
        box.line.width = Pt(1.5)  # crisper outline (P9)
        box.shadow.inherit = False  # flat, no default shadow (P9)
        tf = box.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.text = label
        run = p.runs[0]
        run.font.size = Pt(20)
        run.font.bold = True
        run.font.color.rgb = WHITE
        if prev_right is not None:
            conn = slide.shapes.add_connector(
                MSO_CONNECTOR.STRAIGHT, prev_right, Emu(centers_y), x, Emu(centers_y)
            )
            conn.line.color.rgb = BRAND_SLATE
            conn.line.width = Pt(2.5)
        # Per-box sub-label as a separate borderless textbox below the box (this is
        # a structural slide, not the demo slide, so authored text is fine). (P9)
        cap = slide.shapes.add_textbox(
            x, Emu(int(top) + int(box_h) + 91440), box_w, Emu(457200)
        )
        ctf = cap.text_frame
        ctf.word_wrap = True
        cp = ctf.paragraphs[0]
        cp.alignment = PP_ALIGN.CENTER
        cp.text = sub
        crun = cp.runs[0]
        crun.font.size = Pt(12)
        crun.font.color.rgb = BRAND_SLATE
        prev_right = int(x) + int(box_w)
        x = Emu(int(x) + int(box_w) + int(gap))
    return slide


def _add_demo_slide(prs):
    """DEMO slide: its only text EQUALS a layout placeholder prompt.

    The engine classifies a slide as a clearable 'demo' region when every text
    run on it equals a layout/master placeholder prompt. We build a Title Only
    slide and set the title to the exact default-template title prompt, so the
    slide's only text is an unedited prompt -> demo-detection fires.
    """
    layout = prs.slide_layouts[LO_TITLE_ONLY]
    slide = prs.slides.add_slide(layout)
    sw, sh = prs.slide_width, prs.slide_height
    # Frame the demo slide with NON-TEXT bands only (decoration carrying no text),
    # so _slide_is_demo still sees the title prompt as the slide's ONLY text. (P5)
    _add_band(slide, Emu(0), Emu(0), sw, Emu(228600), BRAND_NAVY)
    _add_band(slide, Emu(0), Emu(228600), sw, Emu(45720), BRAND_AMBER)
    _add_band(slide, Emu(0), Emu(int(sh) - 228600), sw, Emu(228600), BRAND_NAVY)
    title = slide.placeholders[0]
    # Re-read the live prompt so this stays correct if the template changes.
    layout_title = layout.placeholders[0]
    prompt = (
        (layout_title.text or "").strip()
        if getattr(layout_title, "has_text_frame", False)
        else ""
    )
    title.text_frame.text = prompt or TITLE_PROMPT
    return slide


def _add_closing(prs):
    """Closing slide on ``Section Header`` (a distinct closing layout use)."""
    layout = prs.slide_layouts[LO_SECTION]
    slide = prs.slides.add_slide(layout)
    ph = {p.placeholder_format.idx: p for p in slide.placeholders}
    # Re-stack the SAME two Section Header placeholders into a clean title -> body
    # order (override the overlapping default layout coords). (P1)
    ph[0].left, ph[0].top, ph[0].width, ph[0].height = (
        Emu(685800),
        Emu(2057400),
        Emu(7772400),
        Emu(1143000),
    )
    ph[1].left, ph[1].top, ph[1].width, ph[1].height = (
        Emu(685800),
        Emu(3429000),
        Emu(7772400),
        Emu(2286000),
    )
    # Left amber accent bar (decoration only, no text). (P2)
    _add_band(slide, Emu(457200), Emu(2057400), Emu(91440), Emu(3429000), BRAND_AMBER)
    _set_text(ph[0], "Thank You", size=40, bold=True, color=BRAND_NAVY)
    # Anchor the body to the TOP of its frame so the lines sit right below the
    # title instead of dropping to the bottom of the tall body box. (residual fix)
    ph[1].text_frame.vertical_anchor = MSO_ANCHOR.TOP
    _bullets(
        ph[1],
        [
            ("Questions? ", 0),
            ("BrandDocs Corp - synthetic sample deck.", 0),
        ],
        size=18,
    )
    # Make the contact address a REAL mailto hyperlink (PPTX-A5): a run-level
    # a:hlinkClick + slide relationship. Appended to the first bullet so the line
    # reads "Questions? hello@branddocs.example" with the address clickable.
    contact_para = ph[1].text_frame.paragraphs[0]
    mail_run = contact_para.add_run()
    mail_run.text = "hello@branddocs.example"
    mail_run.font.size = Pt(18)
    mail_run.hyperlink.address = "mailto:hello@branddocs.example"
    _set_notes(
        slide,
        "Speaker notes: close on the contact line; the address is a live mailto "
        "link, so a reader can reach the brand team straight from the deck.",
    )
    return slide


def _add_comparison(prs):
    """Comparison-layout slide (PPTX-A9): fills a DISTINCT stock layout's
    placeholders with a side-by-side 'Before vs After brand consistency' story.

    The default template ships a ``Comparison`` layout (TITLE + two BODY headers
    + two OBJECT bodies) that the deck did not previously use; filling its
    placeholders surfaces a body/region slide on a real, previously-unused layout
    (richer region inventory + a content candidate on a distinct layout). Static
    placeholder text, deterministic.
    """
    layout = prs.slide_layouts[LO_COMPARISON]
    slide = prs.slides.add_slide(layout)
    ph = {p.placeholder_format.idx: p for p in slide.placeholders}
    _set_text(
        ph[0], "Before vs After Brand Consistency", size=28, bold=True, color=BRAND_NAVY
    )
    # Left column header + body (idx 1 = BODY header, idx 2 = OBJECT body).
    if 1 in ph:
        _set_text(ph[1], "Before", size=18, bold=True, color=BRAND_SLATE)
    if 2 in ph:
        _bullets(
            ph[2],
            [
                ("Every deck redrew the palette by hand.", 0),
                ("Fonts drifted slide to slide.", 0),
                ("No single source of brand truth.", 0),
            ],
            size=15,
        )
    # Right column header + body (idx 3 = BODY header, idx 4 = OBJECT body).
    if 3 in ph:
        _set_text(ph[3], "After", size=18, bold=True, color=BRAND_TEAL)
    if 4 in ph:
        _bullets(
            ph[4],
            [
                ("One BrandDocs theme drives every surface.", 0),
                ("Arial headings, Calibri body, applied once.", 0),
                ("Outputs match the template by construction.", 0),
            ],
            size=15,
        )
    _set_notes(
        slide,
        "Speaker notes: contrast the two columns; the After column is the payoff "
        "of authoring the palette and fonts at the theme level.",
    )
    return slide


# ---------------------------------------------------------------------------
# Theme1.xml rewrite (clrScheme + fontScheme) via raw lxml (PPTX-A1 / PPTX-A2).
# ---------------------------------------------------------------------------
def _theme_part(prs):
    """Return the package's ``ppt/theme/theme1.xml`` part (the master's theme)."""
    for part in prs.part.package.iter_parts():
        if part.partname.endswith("theme1.xml"):
            return part
    return None


def _rebrand_theme(prs) -> None:
    """Rewrite the EXISTING ``theme1.xml`` clrScheme + fontScheme in place (A1/A2).

    The deck ships with the stock Office theme. This mutates the already-present
    theme part so the deck-wide palette becomes the BrandDocs palette at the
    theme level (the single biggest pptx brand-strength lever the extractor
    reads via ``parse_theme_colors``) and the heading/body font pairing becomes
    Arial/Calibri (read by the theme-fonts path). The edit is a static hex /
    typeface rewrite of an existing part - constants only, no wall-clock - so
    ``freeze_ooxml`` keeps the build byte-reproducible.

    Every ``a:clrScheme`` slot is normalized to an explicit ``a:srgbClr`` (the
    stock ``dk1``/``lt1`` ship as ``a:sysClr``); the extractor resolves both, but
    an explicit srgbClr makes the BrandDocs hex unambiguous in the surfaced
    ``theme.colors``.
    """
    part = _theme_part(prs)
    if part is None:  # pragma: no cover - default template always carries a theme
        return
    root = etree.fromstring(part.blob)
    a = lambda tag: f"{{{A_NS}}}{tag}"  # noqa: E731 - tiny local qname helper

    # --- clrScheme: replace each slot's child with an explicit a:srgbClr. ------
    scheme = root.find(f".//{a('clrScheme')}")
    if scheme is not None:
        for slot, hexval in BRAND_CLRSCHEME.items():
            node = scheme.find(a(slot))
            if node is None:
                continue
            for child in list(node):
                node.remove(child)
            srgb = etree.SubElement(node, a("srgbClr"))
            srgb.set("val", hexval)

    # --- fontScheme: set the major (heading) + minor (body) latin typefaces. ---
    font_scheme = root.find(f".//{a('fontScheme')}")
    if font_scheme is not None:
        for font_tag, typeface in (
            ("majorFont", BRAND_MAJOR_FONT),
            ("minorFont", BRAND_MINOR_FONT),
        ):
            font_el = font_scheme.find(a(font_tag))
            if font_el is None:
                continue
            latin = font_el.find(a("latin"))
            if latin is None:
                latin = etree.SubElement(font_el, a("latin"))
            latin.set("typeface", typeface)

    part._blob = etree.tostring(root, xml_declaration=True, encoding="UTF-8")


# ---------------------------------------------------------------------------
# Subtle slide transitions on section-opening slides via raw lxml (PPTX-A8).
# ---------------------------------------------------------------------------
def _add_transition(slide) -> None:
    """Inject a subtle fixed-duration fade ``p:transition`` into a slide's XML.

    python-pptx does not model slide transitions, so the element is authored
    directly (mirroring the ``p14:sectionLst`` lxml discipline): a single
    ``<p:transition spd="med"><p:fade/></p:transition>`` child appended to the
    ``<p:sld>`` after the ``<p:cSld>``. The ``p:`` prefix is bound explicitly via
    ``nsmap`` so serialization is byte-stable. Static element, no wall-clock.
    """
    sld = slide._element  # <p:sld>
    if sld.find(qn("p:transition")) is not None:
        return
    transition = etree.SubElement(sld, f"{{{P_NS}}}transition", nsmap={"p": P_NS})
    transition.set("spd", "med")
    etree.SubElement(transition, f"{{{P_NS}}}fade")


# ---------------------------------------------------------------------------
# Section list (p14:sectionLst) injection via raw lxml (python-pptx can't reach).
# ---------------------------------------------------------------------------
def _inject_sections(prs, sections) -> None:
    """Inject a real ``p14:sectionLst`` into ``presentation.xml``.

    ``sections`` is a list of ``(name, [slide_index, ...])``. We read the live
    ``p:sldIdLst`` to map slide indices to their ``r:id``/``id`` and emit a
    ``<p:ext uri="{...}"><p14:sectionLst>...</p14:sectionLst></p:ext>`` block
    under ``p:extLst`` - the exact shape the extractor's ``detect_sections``
    walks. Deterministic section GUIDs keep the output reproducible.
    """
    pres = prs.part._element  # <p:presentation>
    p_ns = "http://schemas.openxmlformats.org/presentationml/2006/main"

    # Map slide index -> the sldId 'id' attribute (PowerPoint section sldId uses it).
    sld_id_lst = pres.find(f"{{{p_ns}}}sldIdLst")
    sld_ids = list(sld_id_lst) if sld_id_lst is not None else []
    index_to_id = {i: sld.get("id") for i, sld in enumerate(sld_ids)}

    # Build (or find) the presentation-level extLst, AFTER sldSz/notesSz per schema
    # order; appending at the end of <p:presentation> is schema-valid for extLst.
    ext_lst = pres.find(f"{{{p_ns}}}extLst")
    if ext_lst is None:
        ext_lst = etree.SubElement(pres, f"{{{p_ns}}}extLst")

    ext = etree.SubElement(ext_lst, f"{{{p_ns}}}ext")
    ext.set("uri", SECTION_EXT_URI)
    # Bind the ``p14`` prefix EXPLICITLY on the sectionLst element via nsmap so
    # serialization is deterministic. Relying on lxml's global
    # ``register_namespace`` registry yields a process-dependent prefix
    # (``ns0`` vs ``p14``), which would break byte-reproducibility.
    section_lst = etree.SubElement(
        ext, f"{{{SECTION_NS}}}sectionLst", nsmap={"p14": SECTION_NS}
    )
    # Pin a deterministic GUID per section (index-derived) for reproducibility.
    for si, (name, slide_indices) in enumerate(sections):
        section = etree.SubElement(section_lst, f"{{{SECTION_NS}}}section")
        section.set("name", name)
        guid = "{%08X-0000-4000-8000-%012X}" % (si + 1, si + 1)
        section.set("id", guid)
        sld_id_lst_el = etree.SubElement(section, f"{{{SECTION_NS}}}sldIdLst")
        for idx in slide_indices:
            sid = index_to_id.get(idx)
            if sid is None:
                continue
            sld = etree.SubElement(sld_id_lst_el, f"{{{SECTION_NS}}}sldId")
            sld.set("id", sid)


# ---------------------------------------------------------------------------
# Build.
# ---------------------------------------------------------------------------
def build(out: Path = OUT) -> Path:
    out = Path(out)
    # The shared text-only BrandDocs wordmark. Bytes are computed
    # deterministically by _brandlib, so every example template embeds the SAME
    # logo and re-runs stay byte-identical.
    png_bytes = branddocs_mark_png(640, 160)

    # Materialize the wordmark PNG to a temp file (python-pptx wants a path or a
    # stream; we use a BytesIO-backed temp via NamedTemporaryFile-free path).
    import tempfile

    tmp_png = Path(tempfile.gettempdir()) / "branddocs_template_wordmark.png"
    tmp_png.write_bytes(png_bytes)

    prs = Presentation()  # default template (master + 11 layouts + theme)
    # Drop nothing: the default presentation ships with zero slides.

    # The three named sections and the agenda must agree, so define names first.
    section_names = ["Overview", "Financials", "Closing"]

    # --- Build slides in deck order; collect their 0-based indices. -----------
    s0 = _add_cover(prs, tmp_png)  # 0 cover
    _add_agenda(prs, section_names)  # 1 agenda / section list
    s2 = _add_content_text(prs)  # 2 content-text
    s3 = _add_kpi_dashboard(prs)  # 3 KPI dashboard + native line chart
    s4 = _add_content_table(prs)  # 4 content-table (merged-cell native table)
    s5 = _add_chart(prs)  # 5 native clustered-bar chart
    s6 = _add_revenue_mix(prs)  # 6 native doughnut chart (3rd family)
    s7 = _add_risk_heatmap(prs)  # 7 native heatmap table
    s8 = _add_picture_slide(prs, tmp_png)  # 8 picture
    s9 = _add_comparison(prs)  # 9 Comparison-layout slide
    s10 = _add_smartart_approx(prs)  # 10 grouped-shape "SmartArt"
    s11 = _add_demo_slide(prs)  # 11 DEMO (prompt-only text)
    s12 = _add_closing(prs)  # 12 closing
    _ = (s0, s2, s3, s4, s5, s6, s7, s8, s9, s10, s11, s12)

    # --- Brand the deck theme (clrScheme + fontScheme) in place (A1/A2). -------
    _rebrand_theme(prs)

    # --- Subtle fade transitions on the section-opening slides (A8). ----------
    # The first slide of each named section gets a transition so the deck has a
    # finished cadence; the demo slide (11) is intentionally left untouched.
    for opener in (s0, s4, s8):
        _add_transition(opener)

    # --- Inject the real PowerPoint section list (lxml). ----------------------
    # Section spans (0-based slide indices):
    #   Overview   -> cover, agenda, exec-summary, KPI dashboard (0,1,2,3)
    #   Financials -> table, bar chart, doughnut, heatmap         (4,5,6,7)
    #   Closing    -> picture, comparison, smartart, demo, thanks (8,9,10,11,12)
    sections = [
        ("Overview", [0, 1, 2, 3]),
        ("Financials", [4, 5, 6, 7]),
        ("Closing", [8, 9, 10, 11, 12]),
    ]
    _inject_sections(prs, sections)

    # --- Pin core-properties timestamps for byte-reproducibility. -------------
    fixed = datetime(2026, 1, 15, 0, 0, 0, tzinfo=timezone.utc)
    cp = prs.core_properties
    cp.author = "BrandDocs Corp (synthetic fixture)"
    cp.title = "BrandDocs Corp Quarterly Business Review"
    cp.subject = "Synthetic test fixture - no proprietary content"
    cp.created = fixed
    cp.modified = fixed
    cp.last_modified_by = "BrandDocs Corp (synthetic fixture)"
    cp.revision = 1

    out.parent.mkdir(parents=True, exist_ok=True)
    # Save via a buffer first so we never leave a half-written file on error.
    buf = BytesIO()
    prs.save(buf)
    out.write_bytes(buf.getvalue())
    freeze_ooxml(out)
    return out


if __name__ == "__main__":
    path = build()
    print(f"built {path}")
