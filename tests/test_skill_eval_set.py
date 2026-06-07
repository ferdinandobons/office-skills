# SPDX-License-Identifier: MIT
from __future__ import annotations

import json
import os
import sys
import tempfile
import unittest
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parents[1] / "scripts"))

from docx import Document
from openpyxl import load_workbook
from pptx import Presentation

from brandkit.cli import main

ROOT = Path(__file__).resolve().parents[1]
EVAL_SET = ROOT / "evals" / "skill_eval_set.json"


def _load_eval_set() -> dict:
    return json.loads(EVAL_SET.read_text(encoding="utf-8"))


class SkillEvalSetTest(unittest.TestCase):
    """Run the skill regression prompts against the real example templates."""

    def test_skill_eval_cases_generate_expected_outputs(self) -> None:
        data = _load_eval_set()
        self.assertEqual(data["schema_version"], "branddocs-skill-eval-1")
        self.assertGreaterEqual(len(data["cases"]), 3)
        for case in data["cases"]:
            with self.subTest(case=case["id"]):
                self._run_case(case)

    def _run_case(self, case: dict) -> None:
        with tempfile.TemporaryDirectory() as td:
            tmp = Path(td)
            old_cwd = Path.cwd()
            os.chdir(tmp)
            try:
                template = ROOT / case["template"]
                self.assertTrue(template.is_file(), template)
                self.assertEqual(
                    main(
                        [
                            "extract",
                            "--name",
                            case["brand"],
                            "--template",
                            str(template),
                            "--scope",
                            "project",
                        ]
                    ),
                    0,
                )
                self.assertEqual(
                    main(
                        [
                            "verify",
                            "--name",
                            case["brand"],
                            "--scope",
                            "project",
                            "--qa",
                            "fast",
                        ]
                    ),
                    0,
                )
                if case.get("comprehension"):
                    comprehension_path = tmp / f"{case['id']}.comprehension.json"
                    comprehension_path.write_text(
                        json.dumps(case["comprehension"]),
                        encoding="utf-8",
                    )
                    self.assertEqual(
                        main(
                            [
                                "comprehend",
                                "--name",
                                case["brand"],
                                "--input",
                                str(comprehension_path),
                                "--scope",
                                "project",
                            ]
                        ),
                        0,
                    )
                input_path = tmp / f"{case['id']}.json"
                input_path.write_text(json.dumps(case["input"]), encoding="utf-8")
                output = tmp / f"{case['id']}.{case['kind']}"
                self.assertEqual(
                    main(
                        [
                            "generate",
                            "--name",
                            case["brand"],
                            "--input",
                            str(input_path),
                            "--output",
                            str(output),
                            "--scope",
                            "project",
                            "--qa",
                            "fast",
                        ]
                    ),
                    0,
                )
                self.assertTrue(output.is_file())
                if case["kind"] == "docx":
                    self._assert_docx(output, case["expect"])
                elif case["kind"] == "pptx":
                    self._assert_pptx(output, case["expect"])
                elif case["kind"] == "xlsx":
                    self._assert_xlsx(output, case["expect"])
                else:
                    self.fail(f"unknown eval kind {case['kind']!r}")
            finally:
                os.chdir(old_cwd)

    def _assert_docx(self, output: Path, expect: dict) -> None:
        doc = Document(output)
        text = "\n".join(
            t.text for t in doc.element.iter() if t.tag.endswith("}t") and t.text
        )
        for required in expect.get("required_text") or []:
            self.assertIn(required, text)
        for forbidden in expect.get("forbidden_text") or []:
            self.assertNotIn(forbidden, text)
        self.assertGreaterEqual(len(doc.tables), int(expect.get("min_tables", 0)))

    def _assert_pptx(self, output: Path, expect: dict) -> None:
        prs = Presentation(output)
        texts: list[str] = []
        native_tables = 0
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    texts.append(shape.text)
                if getattr(shape, "has_table", False):
                    native_tables += 1
                    for row in shape.table.rows:
                        texts.extend(cell.text for cell in row.cells)
        text = "\n".join(texts)
        for required in expect.get("required_text") or []:
            self.assertIn(required, text)
        for forbidden in expect.get("forbidden_text") or []:
            self.assertNotIn(forbidden, text)
        self.assertGreaterEqual(native_tables, int(expect.get("min_native_tables", 0)))

    def _assert_xlsx(self, output: Path, expect: dict) -> None:
        wb = load_workbook(output, data_only=False)
        for name, expected in (expect.get("named_values") or {}).items():
            sheet, coord = _first_defined_destination(wb, name)
            self.assertEqual(wb[sheet][coord].value, expected)
        formula_count = 0
        for ws in wb.worksheets:
            for cell in ws._cells.values():
                if isinstance(cell.value, str) and cell.value.startswith("="):
                    formula_count += 1
        self.assertGreaterEqual(formula_count, int(expect.get("min_formula_count", 0)))


def _first_defined_destination(wb, name: str) -> tuple[str, str]:
    defined = wb.defined_names[name]
    destinations = list(defined.destinations)
    if not destinations:
        raise AssertionError(f"defined name {name!r} has no destinations")
    return destinations[0]
